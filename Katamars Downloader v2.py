import requests
from io import StringIO
from html.parser import HTMLParser
import re
import json
import datetime
import tkinter as tk
from tkinter import ttk
from docx import Document
from docx.shared import RGBColor
import os, sys
import pptx
from pptx.dml.color import RGBColor as rgb


class MLStripper(HTMLParser):
    def __init__(self):
        super().__init__()
        self.reset()
        self.strict = False
        self.convert_charrefs= True
        self.text = StringIO()
    def handle_data(self, d):
        self.text.write(d)
    def get_data(self):
        return self.text.getvalue()

def strip_tags(html):
    s = MLStripper()
    s.feed(html)
    return s.get_data()

days_ar= ["الاثنين", "الثلاثاء", "الأربعاء", "الخميس", "الجمعة", "السبت", "الأحد"]

date = datetime.datetime.now()
day_now = date.day
mon_now = date.month
year_now = date.year
day_week_now = days_ar[date.weekday()]
choosed_d=day_now
choosed_m=mon_now
choosed_y=year_now


# print(days_ar[day_week_now])
# nnnn= datetime.datetime(2022, 5, 10, 8, 42, 10)


def main_function(day, mon, year):
    
    docx = Document()
    ppt= pptx.Presentation("stock.pptx")
    # ppt= pptx.Presentation("23 ابيب.pptx")
    title_slide_layout = ppt.slide_layouts[2]
    


    stat_variable.set("جارى التحميل\t"+"0%")
    r.update()
    x = requests.get(f'http://katamars.avabishoy.com/api/Katamars/GetReadings?day={day}&katamrsSourceId=1&month={mon}&synaxariumSourceId=1&year={year}')
    text= x.content.decode()

    newtext = json.loads(text)

    polis = newtext["polis"]
    apraksees = newtext["apraksees"]
    kathilycon = newtext["kathilycon"]
    gospel = newtext["gospel"]
    sneKsar = newtext["synaxarium"]

    stat_variable.set("10%")
    r.update()
    
    txt_polis = strip_tags(polis)
    txt_polis = re.sub(r'\(.*?\)', "\n", txt_polis)
    txt_polis = re.sub(r'\t', "", txt_polis)
    txt_polis = txt_polis.replace(":","\n")

    stat_variable.set("20%")
    r.update()

    txt_apraksees = strip_tags(apraksees)
    txt_apraksees = re.sub(r'\(.*?\)', "\n", txt_apraksees)
    txt_apraksees = re.sub(r'\t', "", txt_apraksees)
    txt_apraksees = txt_apraksees.replace(":","\n")

    stat_variable.set("25%")
    r.update()

    txt_kathilycon = strip_tags(kathilycon)
    txt_kathilycon = re.sub(r'\(.*?\)', "\n", txt_kathilycon)
    txt_kathilycon = re.sub(r'\t', "", txt_kathilycon)
    txt_kathilycon = txt_kathilycon.replace(":","\n")

    stat_variable.set("30%")
    r.update()

    txt_gospel = strip_tags(gospel)
    txt_gospel = re.sub(r'\(.*?\)', "\n",  txt_gospel)
    txt_gospel = re.sub(r'\t', "", txt_gospel)
    txt_gospel = txt_gospel.replace("هللويا", "\nهللويا")


    stat_variable.set("37%")
    r.update()
    txt_sneKsar=""
    for i in sneKsar:
        txt_sneKsar = txt_sneKsar+i["title"]+ "\n"

    stat_variable.set("40%")
    r.update()


    coptic_months=["توت", "بابه", "هاتور", "كهيك", "طوبة", "أمشير", "برمهات", "برمودة", "بشنس", "بؤونة", "أبيب", "مسرى»"]

    x = requests.get(f'http://katamars.avabishoy.com/api/Katamars/GetCopticDate?day={day}&month={mon}&year={year}')
    coptic_date = x.content.decode()
    coptic_date = json.loads(coptic_date)
    coptic_mon = coptic_date["month"]
    coptic_mon= coptic_months[coptic_mon-1]

    stat_variable.set("50%")
    r.update()



    final_txt= f"قطمارس ليوم {youm_variable.get()} {day}/{mon}/{year} و يوافق قبطيا  {coptic_date['day']}-{coptic_mon}-{coptic_date['year']}" + "\n\n" + txt_polis  + "\n\n" + txt_kathilycon + "\n\n" + txt_apraksees + "\nالسنكسار\n" + f"اليوم {coptic_date['day']} من الشهر المبارك {coptic_mon}" + txt_sneKsar + "\n الانجيل \nا" + txt_gospel 

    if selected_txt.get() :
        with open("Katamars.txt", "w") as text_file:
            text_file.write(final_txt)

    stat_variable.set("57%")
    r.update()

    if selected_word.get() :

        list_polis= txt_polis.splitlines()
        list_kathilycon= txt_kathilycon.splitlines()
        list_apraksees= txt_apraksees.splitlines()
        list_gospel= txt_gospel.splitlines()
        list_sneKsar= txt_sneKsar.splitlines()
        line_word_count= 10

        docx.add_paragraph(f"قطمارس ليوم {youm_variable.get()} {day}/{mon}/{year} و يوافق قبطيا  {coptic_date['day']}-{coptic_mon}-{coptic_date['year']}")
        if selected_ppt.get():
            slide = ppt.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = f"قطمارس ليوم {youm_variable.get()} {day}/{mon}/{year} و يوافق قبطيا  {coptic_date['day']}-{coptic_mon}-{coptic_date['year']}"
#-------------------------   word   البولس    -----------------------------------------#
        
        if selected_ppt.get() and selected_word.get():
            print(1)
            process_list_word_ppt(list_polis, docx, ppt, title_slide_layout, line_word_count)
        
        elif selected_word.get() :
            print(2)
            process_list_wordonly(list_polis, docx, ppt, title_slide_layout, line_word_count)

        elif selected_ppt.get() :
            print(3)
            process_list_pptonly(list_polis, docx, ppt, title_slide_layout, line_word_count)
        # for i in list_polis :
        #     if i == list_polis[0] : 
        #         run = docx.add_paragraph().add_run(i)
        #         font = run.font
        #         font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        #         if selected_ppt.get():
        #             slide = ppt.slides.add_slide(title_slide_layout)
        #             title = slide.shapes.title
        #             title.text_frame.text = i
        #             font = title.text_frame.paragraphs[0].runs[0].font
        #             font.bold = True
        #             font.color.rgb = rgb(0xFF, 0x00, 0x00)

        #     elif i == list_polis[1]: 
        #         docx.add_paragraph(i)
        #         if selected_ppt.get():
        #             slide = ppt.slides.add_slide(title_slide_layout)
        #             title = slide.shapes.title
        #             title.text = i
        #     else:
        #         words= i.split()
        #         line=""
        #         if len(words) > line_word_count :
        #             for word in words :
        #                 line = line + word + " "
        #                 if len(line.split()) >= line_word_count: 
        #                     docx.add_paragraph(line)
        #                     if selected_ppt.get():
        #                         slide = ppt.slides.add_slide(title_slide_layout)
        #                         title = slide.shapes.title
        #                         title.text = line
        #                     line=""
        #             dd= len(words)%10
        #             if dd != 0 :
        #                 line= " ".join(words[-dd:])
        #                 docx.add_paragraph(line)
        #                 if selected_ppt.get():
        #                     slide = ppt.slides.add_slide(title_slide_layout)
        #                     title = slide.shapes.title
        #                     title.text = line
                    
        #         else : 
        #             docx.add_paragraph(i)  
        #             if selected_ppt.get():
        #                 slide = ppt.slides.add_slide(title_slide_layout)
        #                 title = slide.shapes.title
        #                 title.text = i

        stat_variable.set("72%")
        r.update()
        ppt.save("Katamars.pptx")

#-------------------------   word   كاثوليكون    -----------------------------------------#


        for i in list_kathilycon :
            if i == list_kathilycon[0] : 
                run = docx.add_paragraph().add_run(i)
                font = run.font
                font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

            elif i == list_kathilycon[1]: 
                docx.add_paragraph(i)
            
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word +" "
                        if len(line.split()) >= line_word_count: 
                            docx.add_paragraph(line)
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        docx.add_paragraph(line)

                else : 
                    docx.add_paragraph(i) 

#-------------------------   word   الابركسيس    -----------------------------------------#


        for i in list_apraksees :
            if i == list_apraksees[0] : 
                run = docx.add_paragraph().add_run(i)
                font = run.font
                font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

            elif i == list_apraksees[1]: 
                docx.add_paragraph(i)
            
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word +" "
                        if len(line.split()) >= line_word_count: 
                            docx.add_paragraph(line)
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        docx.add_paragraph(line)
                else : 
                    docx.add_paragraph(i) 

        stat_variable.set("88%")
        r.update()
        
#-------------------------   word   السنكسار    -----------------------------------------#


        run = docx.add_paragraph().add_run("\nالسنكسار\n")
        font = run.font
        font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        
        docx.add_paragraph(f"اليوم {coptic_date['day']} من الشهر المبارك {coptic_mon}")
        for i in list_sneKsar:
            docx.add_paragraph(i)

#-------------------------   word   الانجيل    -----------------------------------------#


        run = docx.add_paragraph().add_run("\n الانجيل \n")
        font = run.font
        font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

        for i in list_gospel :
            if i == list_gospel[0] or  i.__contains__("من إنجيل"): 
                run = docx.add_paragraph().add_run(i)
                font = run.font
                font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word +" "
                        if len(line.split()) >= line_word_count: 
                            docx.add_paragraph(line)
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        docx.add_paragraph(line)
                else : 
                    docx.add_paragraph(i) 

        docx.save('Katamars.docx')

    if selected_ppt.get(): 
        print("hi")

    stat_variable.set( "تمام\t" + "100%")
    r.update()

def process_list_word_ppt(in_list, docx, ppt, title_slide_layout, line_word_count) :
    for i in in_list :
            if i == in_list[0] : 
                run = docx.add_paragraph().add_run(i)
                font = run.font
                font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
                if selected_ppt.get():
                    slide = ppt.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text_frame.text = i
                    font = title.text_frame.paragraphs[0].runs[0].font
                    font.bold = True
                    font.color.rgb = rgb(0xFF, 0x00, 0x00)

            elif i == in_list[1]: 
                docx.add_paragraph(i)
                if selected_ppt.get():
                    slide = ppt.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = i
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word + " "
                        if len(line.split()) >= line_word_count: 
                            docx.add_paragraph(line)
                            if selected_ppt.get():
                                slide = ppt.slides.add_slide(title_slide_layout)
                                title = slide.shapes.title
                                title.text = line
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        docx.add_paragraph(line)
                        if selected_ppt.get():
                            slide = ppt.slides.add_slide(title_slide_layout)
                            title = slide.shapes.title
                            title.text = line
                    
                else : 
                    docx.add_paragraph(i)  
                    if selected_ppt.get():
                        slide = ppt.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        title.text = i

def process_list_wordonly(in_list, docx, ppt, title_slide_layout, line_word_count) :
    for i in in_list :
            if i == in_list[0] : 
                run = docx.add_paragraph().add_run(i)
                font = run.font
                font.color.rgb = RGBColor(0xFF, 0x00, 0x00)

            elif i == in_list[1]: 
                docx.add_paragraph(i)
                
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word + " "
                        if len(line.split()) >= line_word_count: 
                            docx.add_paragraph(line)
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        docx.add_paragraph(line)
                else : 
                    docx.add_paragraph(i)  

def process_list_pptonly(in_list, docx, ppt, title_slide_layout, line_word_count) :
    for i in in_list :
            if i == in_list[0] : 
                slide = ppt.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text_frame.text = i
                font = title.text_frame.paragraphs[0].runs[0].font
                font.bold = True
                font.color.rgb = rgb(0xFF, 0x00, 0x00)

            elif i == in_list[1]: 
                slide = ppt.slides.add_slide(title_slide_layout)
                title = slide.shapes.title
                title.text = i
            else:
                words= i.split()
                line=""
                if len(words) > line_word_count :
                    for word in words :
                        line = line + word + " "
                        if len(line.split()) >= line_word_count: 
                            slide = ppt.slides.add_slide(title_slide_layout)
                            title = slide.shapes.title
                            title.text = line
                            line=""
                    dd= len(words)%10
                    if dd != 0 :
                        line= " ".join(words[-dd:])
                        slide = ppt.slides.add_slide(title_slide_layout)
                        title = slide.shapes.title
                        title.text = line
                
                else : 
                    slide = ppt.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    title.text = i

                    
def choose_day(d, m, y):
    global choosed_d, choosed_m, choosed_y 
    delta_d= d-choosed_d
    delta_m= m-choosed_m
    delta_y= y-choosed_y
    delay= delta_d + delta_m*30 + delta_y*365
    if (delay > 0):
        x=datetime.datetime(choosed_y, choosed_m, choosed_d) + datetime.timedelta(days=delay)
    else :
        x=datetime.datetime(choosed_y, choosed_m, choosed_d) - datetime.timedelta(days=0-delay)
    choosed_d = x.day
    choosed_m = x.month
    choosed_y = x.year
    newyoum=days_ar[x.weekday()]
    youm_variable.set(newyoum)
    years_variable.set(years_op[choosed_y-(year_now-4)]) 
    months_variable.set(months_op[choosed_m]) 
    days_variable.set(days_op[choosed_d]) 
    stat_variable.set("جاهز")

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)

r = tk.Tk()
r.title('Katamars')
r.geometry("800x250")
r.wm_minsize(width=500, height=300)
r.iconbitmap(resource_path("favicon.ico"))

youm_variable = tk.StringVar()
days_variable = tk.StringVar()
months_variable = tk.StringVar()
years_variable = tk.StringVar()
stat_variable = tk.StringVar()
selected_word = tk.StringVar()
selected_txt = tk.StringVar()
selected_ppt = tk.StringVar()


style = ttk.Style()
style.configure("BW.TLabel", foreground="black", background="white", font=('Arial', 14))
style.configure("stat.TLabel", foreground="white", background="#007ACC", font=('Arial', 11,"bold"))
style.configure('my.TButton', font=('Arial', 14))
style.configure('mooo.TButton', font=('Sans', 14, "bold"))
style.configure('my.TMenubutton', background="white")
style.configure('my.TCheckbutton', background="white")




date= tk.Label(r, text=f"تاريخ اليوم : {day_now}/{mon_now}/{year_now}", anchor="e", font= ("Arial, 15"), padx=10, pady=10)
date.pack(fill="x")


############################################################ FRAME ################################################################
frame0= tk.Frame(r, pady=10, background="white")
frame0.pack(fill="x")
frame_date_pick= tk.Frame(frame0, pady=10, background="white")
frame_date_pick.pack()
frame_output= tk.Frame(frame0, pady=10, background="white")
frame_output.pack()
 

button1 = ttk.Button(frame_date_pick, text='التالى', style="my.TButton", command=lambda: choose_day(choosed_d+1, choosed_m, choosed_y) )
button1.grid(row=0, column=1, padx= 10)

days_op = [i for i in range(0, 30+2)]
months_op = [i for i in range(0, 12+1)]
years_op = [i for i in range(year_now-4, year_now+3)]

menu_y = ttk.OptionMenu(frame_date_pick, years_variable, *years_op, style="my.TMenubutton", command=lambda yy: choose_day(choosed_d, choosed_m, yy))
years_variable.set(years_op[choosed_y-(year_now-4)]) 
menu_y.grid(row=0, column=2)

menu_m = ttk.OptionMenu(frame_date_pick, months_variable, *months_op, style="my.TMenubutton", command=lambda mm: choose_day(choosed_d, mm, choosed_y))
months_variable.set(months_op[choosed_m]) 
menu_m.grid(row=0, column=3)

menu_d = ttk.OptionMenu(frame_date_pick, days_variable, *days_op, style="my.TMenubutton", command=lambda dd: choose_day(dd, choosed_m, choosed_y))
days_variable.set(days_op[choosed_d]) 
menu_d.grid(row=0, column=4)


youm= ttk.Label(frame_date_pick, textvariable=youm_variable, style="BW.TLabel")
youm_variable.set(day_week_now)
youm.grid(row=0, column=5, padx=10)

button2 = ttk.Button(frame_date_pick, text='السابق', style="my.TButton", command=lambda: choose_day(choosed_d-1, choosed_m, choosed_y))
button2.grid(row=0, column=6, padx= 10)

l= tk.Label(frame_date_pick, text=" :اختار اليوم", background="white")
l.grid(row=0, column=7)

output= tk.Label(frame_output, text=" :اختار الملف الخارج", background="white")
output.grid(row=1, column=7, pady=10, padx= 10)

opt_wrd= ttk.Checkbutton(frame_output, text="ملف وورد " , style="my.TCheckbutton", variable=selected_word)
opt_wrd.grid(row=1, column=6, padx=20)

opt_txt= ttk.Checkbutton(frame_output, text=" TXT ملف" , style="my.TCheckbutton", variable=selected_txt)
opt_txt.grid(row=1, column=5, padx=20)

opt_ppt= ttk.Checkbutton(frame_output, text=" ملف باور بوينت" , style="my.TCheckbutton", variable=selected_ppt)
opt_ppt.grid(row=1, column=4, padx=20)

selected_txt.set(1)
selected_word.set(1)
selected_ppt.set(1)


############################################################ FRAME ################################################################

button3 = ttk.Button(r, text='تحميل القطمارس', width=25, style="mooo.TButton", command=lambda : main_function(choosed_d, choosed_m, choosed_y) , padding=(10,10,10,10))
button3.pack(pady=20)

stat= ttk.Label(r, textvariable=stat_variable, style="stat.TLabel", anchor="e")
stat_variable.set("جاهز")
stat.pack(fill="x", side="bottom")

r.mainloop()
